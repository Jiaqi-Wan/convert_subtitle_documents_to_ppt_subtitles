from docx import Document
from pptx import Presentation
import re


def word_to_ppt_with_template(word_path, template_path, output_path, delimiter='/'):
    """核心转换函数

    参数：
    word_path: Word文档路径(.docx)
    template_path: 模板PPT路径(.pptx)
    output_path: 输出PPT路径
    delimiter: 句子分隔符
    """

    # 读取Word内容
    doc = Document(word_path)
    full_text = '\n'.join([para.text for para in doc.paragraphs if para.text])
    sentences = re.split(r'\s*' + re.escape(delimiter) + r'\s*', full_text)
    sentences = [s.strip() for s in sentences if s.strip()]

    # 打开模板文件
    prs = Presentation(template_path)

    # 获取模板中的版式（假设使用第一个版式）
    slide_layout = prs.slide_layouts[0]  # 根据实际模板调整索引

    # 查找模板中的占位符
    for shape in slide_layout.placeholders:
        if shape.is_placeholder and shape.placeholder_format.type == 1:  # 1表示标题占位符
            text_frame = shape.text_frame
            break

    # 批量生成幻灯片
    for sentence in sentences:
        new_slide = prs.slides.add_slide(slide_layout)

        # 应用模板样式到新幻灯片
        for shape in new_slide.shapes:
            if shape.is_placeholder:
                # 清空原有内容但保留样式
                if shape.has_text_frame:
                    shape.text_frame.clear()
                    # 添加新文本（自动继承模板样式）
                    p = shape.text_frame.paragraphs[0]
                    p.text = sentence
                    # 保持段落样式不变
                    p.level = 0

    prs.save(output_path)


if __name__ == "__main__":
    # 配置区 ==================================
    word_path = "C:/Users/Kingway/Desktop/2025宝马集团创新日活动方案1.docx"  # 输入Word文件
    template_path = "D:/字幕测试.pptx"  # 模板文件
    output_path = "D:/字幕测试.pptx"  # 输出文件
    delimiter = "/"  # 分隔符
    # =======================================

    word_to_ppt_with_template(word_path, template_path, output_path, delimiter)
    print(f"Success！共生成 {len(Document(word_path).paragraphs)} 页")
